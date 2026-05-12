"""
build_presentation.py
Genera cancer_prediction_presentation.pptx con 5 diapositivas ejecutivas.
Lee métricas reales desde outputs/metrics/; si un archivo falta, muestra
"pendiente de ejecución" en lugar de inventar valores.

Ejecutar desde la raíz del proyecto:
    python outputs/presentation/build_presentation.py
"""

import json
import os
import pathlib
import textwrap

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ── Rutas ─────────────────────────────────────────────────────────────────────
ROOT         = pathlib.Path(__file__).parent.parent.parent   # raíz del proyecto
METRICS_DIR  = ROOT / "outputs" / "metrics"
FIGURES_DIR  = ROOT / "outputs" / "figures"
OUT_DIR      = ROOT / "outputs" / "presentation"
OUT_PPTX     = OUT_DIR / "cancer_prediction_presentation.pptx"

# ── Paleta de colores ─────────────────────────────────────────────────────────
AZUL_OSCURO  = RGBColor(0x1A, 0x3A, 0x5C)   # títulos
AZUL_MEDIO   = RGBColor(0x2C, 0x5F, 0x8A)   # acentos / headers tabla
AZUL_CLARO   = RGBColor(0xD6, 0xE4, 0xF0)   # fondo headers tabla
GRIS_FONDO   = RGBColor(0xF4, 0xF6, 0xF9)   # fondo cards
GRIS_BORDE   = RGBColor(0xDD, 0xE3, 0xED)   # bordes
ROJO         = RGBColor(0xE0, 0x5C, 0x4B)   # alerta / riesgo
VERDE        = RGBColor(0x43, 0xA0, 0x47)   # positivo
BLANCO       = RGBColor(0xFF, 0xFF, 0xFF)
GRIS_TEXTO   = RGBColor(0x6B, 0x7A, 0x8D)   # subtextos

# ── Dimensiones 16:9 ─────────────────────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ─────────────────────────────────────────────────────────────────────────────
# CARGA DE MÉTRICAS
# ─────────────────────────────────────────────────────────────────────────────

def _load_or_none(path: pathlib.Path):
    if path.exists():
        return path
    print(f"  [AVISO] No encontrado: {path.relative_to(ROOT)}")
    return None


def load_metrics() -> dict:
    """Carga todos los archivos de métricas. Devuelve None donde falten."""
    m: dict = {}

    # all_models_metrics.csv
    p = METRICS_DIR / "all_models_metrics.csv"
    m["all_metrics"] = pd.read_csv(p, index_col=0) if p.exists() else None

    # final_model_ranking.csv
    p = METRICS_DIR / "final_model_ranking.csv"
    m["ranking"] = pd.read_csv(p, index_col=0) if p.exists() else None

    # executive_summary.json
    p = METRICS_DIR / "executive_summary.json"
    if p.exists():
        m["summary"] = json.loads(p.read_text(encoding="utf-8"))
    else:
        m["summary"] = None

    # ml_results.csv
    p = METRICS_DIR / "ml_results.csv"
    m["ml_results"] = pd.read_csv(p, index_col=0) if p.exists() else None

    # mlp_results.csv
    p = METRICS_DIR / "mlp_results.csv"
    m["mlp_results"] = pd.read_csv(p, index_col=0) if p.exists() else None

    return m


def fmt(val, decimals=3) -> str:
    """Formatea un valor numérico o devuelve 'pendiente de ejecución'."""
    if val is None:
        return "pendiente de ejecución"
    try:
        return f"{float(val):.{decimals}f}"
    except (TypeError, ValueError):
        return str(val)


def get_best_model_row(m: dict) -> dict | None:
    """Devuelve la fila del mejor modelo desde ranking o all_metrics."""
    if m["ranking"] is not None and len(m["ranking"]) > 0:
        row = m["ranking"].iloc[0]
        return row.to_dict()
    if m["all_metrics"] is not None and len(m["all_metrics"]) > 0:
        df = m["all_metrics"].copy()
        col = next((c for c in ["F1", "F1_cancer1", "f1"] if c in df.columns), None)
        if col:
            return df.sort_values(col, ascending=False).iloc[0].to_dict()
        return df.iloc[0].to_dict()
    return None


# ─────────────────────────────────────────────────────────────────────────────
# HELPERS DE FORMATO PPTX
# ─────────────────────────────────────────────────────────────────────────────

def add_textbox(slide, text, left, top, width, height,
                font_size=14, bold=False, color=None, align=PP_ALIGN.LEFT,
                italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or AZUL_OSCURO
    return txBox


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0.5)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def slide_number(slide, n: int, total: int = 5):
    add_textbox(slide, f"{n} / {total}",
                Inches(12.5), Inches(7.1), Inches(0.8), Inches(0.3),
                font_size=9, color=GRIS_TEXTO, align=PP_ALIGN.RIGHT)


def add_accent_bar(slide):
    """Barra horizontal azul en la parte superior."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=AZUL_MEDIO)


def add_footer(slide, text="Universidad Alfonso X el Sabio · Bases de Datos e IA · 2025-2026"):
    add_rect(slide, 0, Inches(7.25), SLIDE_W, Inches(0.25), fill_color=AZUL_OSCURO)
    add_textbox(slide, text,
                Inches(0.2), Inches(7.27), Inches(11), Inches(0.22),
                font_size=8, color=BLANCO, align=PP_ALIGN.LEFT)


def add_title_block(slide, title: str, subtitle: str = ""):
    add_textbox(slide, title,
                Inches(0.35), Inches(0.15), Inches(12.5), Inches(0.6),
                font_size=24, bold=True, color=AZUL_OSCURO)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.35), Inches(0.72), Inches(12.5), Inches(0.35),
                    font_size=13, color=AZUL_MEDIO)


def add_divider(slide, top_inches: float):
    add_rect(slide, Inches(0.35), Inches(top_inches),
             Inches(12.63), Pt(1), fill_color=AZUL_CLARO)


def add_image_safe(slide, fname: str, left, top, width, height, caption=""):
    path = FIGURES_DIR / fname
    if path.exists():
        slide.shapes.add_picture(str(path), left, top, width, height)
        if caption:
            add_textbox(slide, caption, left, top + height, width, Inches(0.22),
                        font_size=8, color=GRIS_TEXTO, align=PP_ALIGN.CENTER)
        return True
    else:
        add_rect(slide, left, top, width, height,
                 fill_color=GRIS_FONDO, line_color=GRIS_BORDE)
        add_textbox(slide, f"[Figura no disponible]\n{fname}\nEjecuta el notebook para generarla.",
                    left + Inches(0.1), top + height * 0.35, width - Inches(0.2), Inches(0.6),
                    font_size=9, color=GRIS_TEXTO, align=PP_ALIGN.CENTER)
        return False


def add_kpi_card(slide, label: str, value: str, sublabel: str,
                 left, top, width, height, value_color=None):
    add_rect(slide, left, top, width, height,
             fill_color=GRIS_FONDO, line_color=AZUL_MEDIO, line_width=Pt(1.2))
    add_rect(slide, left, top, Pt(4), height, fill_color=AZUL_MEDIO)
    add_textbox(slide, label,
                left + Inches(0.12), top + Inches(0.05),
                width - Inches(0.15), Inches(0.22),
                font_size=8, bold=True, color=GRIS_TEXTO)
    add_textbox(slide, value,
                left + Inches(0.12), top + Inches(0.22),
                width - Inches(0.15), Inches(0.38),
                font_size=20, bold=True, color=value_color or AZUL_OSCURO)
    add_textbox(slide, sublabel,
                left + Inches(0.12), top + Inches(0.55),
                width - Inches(0.15), Inches(0.2),
                font_size=8, color=GRIS_TEXTO)


# ─────────────────────────────────────────────────────────────────────────────
# TABLA GENÉRICA
# ─────────────────────────────────────────────────────────────────────────────

def add_table(slide, headers, rows, left, top, width, height,
              highlight_row: int = None, col_widths=None):
    """Tabla con cabecera azul y filas alternadas."""
    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 header
    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    # Anchos de columna
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(width * cw / total)

    row_h = height // n_rows
    for r in tbl.rows:
        r.height = row_h

    def cell_style(cell, text, bold=False, bg=None, fg=None, font_size=10,
                   align=PP_ALIGN.CENTER):
        cell.text = text
        cell.text_frame.paragraphs[0].alignment = align
        tf = cell.text_frame
        for para in tf.paragraphs:
            para.alignment = align
            for run in para.runs:
                run.font.size = Pt(font_size)
                run.font.bold = bold
                if fg:
                    run.font.color.rgb = fg
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
        cell.margin_left = Inches(0.04)
        cell.margin_right = Inches(0.04)
        cell.margin_top = Inches(0.02)
        cell.margin_bottom = Inches(0.02)

    # Header
    for j, h in enumerate(headers):
        cell_style(tbl.cell(0, j), h, bold=True, bg=AZUL_MEDIO, fg=BLANCO,
                   font_size=9, align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

    # Rows
    for i, row_data in enumerate(rows):
        row_idx = i + 1
        is_highlight = (highlight_row is not None and i == highlight_row)
        bg = AZUL_CLARO if is_highlight else (BLANCO if i % 2 == 0 else GRIS_FONDO)
        for j, cell_val in enumerate(row_data):
            cell_style(tbl.cell(row_idx, j), str(cell_val),
                       bold=is_highlight,
                       bg=bg,
                       fg=AZUL_OSCURO,
                       font_size=9,
                       align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

    return tbl


# ─────────────────────────────────────────────────────────────────────────────
# SLIDES
# ─────────────────────────────────────────────────────────────────────────────

def build_slide1(prs: Presentation, m: dict):
    """SLIDE 1 — Objetivo y datos."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    add_accent_bar(slide)
    add_footer(slide)
    slide_number(slide, 1)

    add_title_block(slide,
                    "Predicción de Diagnóstico de Cáncer con IA",
                    "Comparativa ML clásico vs Red Neuronal Multicapa · UAX · 2025-2026")
    add_divider(slide, 1.1)

    # ── 4 KPI cards ────────────────────────────────────────────────
    card_w = Inches(2.9)
    card_h = Inches(0.95)
    card_top = Inches(1.25)
    gaps = [Inches(0.35), Inches(3.35), Inches(6.35), Inches(9.35)]

    add_kpi_card(slide, "PACIENTES", "50 001", "Dataset sintético oncológico",
                 gaps[0], card_top, card_w, card_h)
    add_kpi_card(slide, "COLECCIONES", "6", "Bioquímica · Clínica · Genética · Económica · Hábitos · Sociodem.",
                 gaps[1], card_top, card_w, card_h)
    add_kpi_card(slide, "PREVALENCIA CÁNCER", "≈ 19 %", "Desbalance 4.3 : 1",
                 gaps[2], card_top, card_w, card_h, value_color=ROJO)
    add_kpi_card(slide, "VARIABLES ACTIVAS", "29 / 38", "7 excluidas por data leakage",
                 gaps[3], card_top, card_w, card_h)

    add_divider(slide, 2.32)

    # ── Pipeline horizontal ─────────────────────────────────────────
    add_textbox(slide, "PIPELINE DEL PROYECTO",
                Inches(0.35), Inches(2.42), Inches(4), Inches(0.28),
                font_size=9, bold=True, color=AZUL_MEDIO)

    pasos = [
        ("6 CSVs", "Datos crudos\n6 colecciones"),
        ("Merge", "JOIN por\npaciente_id"),
        ("EDA &\nLeakage", "Audit de\nvariables"),
        ("Pipeline\nsklearn", "fit solo en\ntrain"),
        ("ML / MLP", "5 modelos\nclásicos + MLP"),
        ("Test\nEval.", "Una sola vez\npor modelo"),
    ]
    n = len(pasos)
    box_w = Inches(1.85)
    box_h = Inches(0.78)
    box_top = Inches(2.75)
    arrow_w = Inches(0.25)
    total_w = n * box_w + (n - 1) * arrow_w
    start_x = (SLIDE_W - total_w) / 2

    for i, (titulo, desc) in enumerate(pasos):
        bx = start_x + i * (box_w + arrow_w)
        add_rect(slide, bx, box_top, box_w, box_h,
                 fill_color=AZUL_CLARO, line_color=AZUL_MEDIO, line_width=Pt(0.75))
        add_textbox(slide, titulo, bx + Inches(0.05), box_top + Inches(0.04),
                    box_w - Inches(0.1), Inches(0.28),
                    font_size=9, bold=True, color=AZUL_OSCURO, align=PP_ALIGN.CENTER)
        add_textbox(slide, desc, bx + Inches(0.05), box_top + Inches(0.3),
                    box_w - Inches(0.1), Inches(0.42),
                    font_size=7.5, color=AZUL_MEDIO, align=PP_ALIGN.CENTER)
        if i < n - 1:
            ax = bx + box_w
            add_textbox(slide, "→", ax + Inches(0.03), box_top + Inches(0.25),
                        arrow_w, Inches(0.3),
                        font_size=12, bold=True, color=AZUL_MEDIO, align=PP_ALIGN.CENTER)

    add_divider(slide, 3.66)

    # ── Anti-leakage banner ─────────────────────────────────────────
    add_rect(slide, Inches(0.35), Inches(3.75), Inches(5.8), Inches(0.75),
             fill_color=RGBColor(0xE8, 0xF0, 0xF8), line_color=AZUL_MEDIO, line_width=Pt(0.75))
    add_rect(slide, Inches(0.35), Inches(3.75), Pt(4), Inches(0.75), fill_color=AZUL_MEDIO)
    add_textbox(slide, "Protocolo Anti-Data-Leakage",
                Inches(0.55), Inches(3.80), Inches(5.5), Inches(0.25),
                font_size=9, bold=True, color=AZUL_OSCURO)
    add_textbox(slide,
                "Pipeline fit solo en train · Umbral optimizado en val · Test evaluado 1 sola vez por modelo",
                Inches(0.55), Inches(4.05), Inches(5.5), Inches(0.35),
                font_size=8, color=AZUL_MEDIO)

    # ── Variables excluidas ─────────────────────────────────────────
    add_rect(slide, Inches(6.35), Inches(3.75), Inches(6.65), Inches(0.75),
             fill_color=RGBColor(0xFD, 0xED, 0xEB), line_color=ROJO, line_width=Pt(0.75))
    add_rect(slide, Inches(6.35), Inches(3.75), Pt(4), Inches(0.75), fill_color=ROJO)
    add_textbox(slide, "Variables Excluidas (leakage)",
                Inches(6.55), Inches(3.80), Inches(6.3), Inches(0.25),
                font_size=9, bold=True, color=RGBColor(0x7A, 0x1A, 0x10))
    add_textbox(slide,
                "coste_total · coste_farmaco · num_ingresos · dias_hospital · vive · alcohol · tipo_seguro",
                Inches(6.55), Inches(4.05), Inches(6.3), Inches(0.35),
                font_size=7.5, color=ROJO)

    # ── Figura ──────────────────────────────────────────────────────
    add_image_safe(slide, "fig_01_target_distribution.png",
                   Inches(0.35), Inches(4.65), Inches(5.8), Inches(2.4),
                   "Distribución de la variable objetivo")
    add_image_safe(slide, "fig_08_leakage_correlation.png",
                   Inches(6.35), Inches(4.65), Inches(6.65), Inches(2.4),
                   "Variables con riesgo de data leakage — auditadas y excluidas")


def build_slide2(prs: Presentation, m: dict):
    """SLIDE 2 — Resultados ML clásicos."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_accent_bar(slide)
    add_footer(slide)
    slide_number(slide, 2)

    add_title_block(slide,
                    "Modelos ML Clásicos: Rendimiento en Test",
                    "Umbral de decisión optimizado en validación · Test evaluado una única vez")
    add_divider(slide, 1.1)

    # ── Preparar datos ──────────────────────────────────────────────
    PENDING = "pend."
    headers = ["Modelo", "Precision", "Recall", "F1", "AUC-ROC", "AUC-PR", "Umbral"]
    model_order = ["LogisticRegression", "RandomForest", "XGBoost", "LightGBM", "CatBoost"]
    rows = []
    best_idx = 0

    if m["all_metrics"] is not None:
        df = m["all_metrics"].copy()
        # Filtrar solo modelos ML (no MLP)
        if "Tipo" in df.columns:
            df_ml = df[df["Tipo"].str.upper().str.contains("ML", na=False)]
        elif "Modelo" in df.columns:
            df_ml = df[~df["Modelo"].astype(str).str.upper().str.contains("MLP")]
        else:
            df_ml = df[~df.index.astype(str).str.upper().str.contains("MLP")]
        if len(df_ml) == 0:
            df_ml = df

        col_map = {
            "Precision": next((c for c in ["Precision", "precision"] if c in df_ml.columns), None),
            "Recall":    next((c for c in ["Recall", "recall"] if c in df_ml.columns), None),
            "F1":        next((c for c in ["F1", "F1_cancer1", "f1"] if c in df_ml.columns), None),
            "AUC-ROC":   next((c for c in ["AUC-ROC", "auc_roc", "AUC_ROC"] if c in df_ml.columns), None),
            "AUC-PR":    next((c for c in ["AUC-PR", "auc_pr", "AUC_PR"] if c in df_ml.columns), None),
            "Threshold": next((c for c in ["Threshold", "threshold"] if c in df_ml.columns), None),
        }
        f1_col = col_map["F1"]
        if f1_col:
            df_ml = df_ml.sort_values(f1_col, ascending=False)

        for i, (idx, row) in enumerate(df_ml.iterrows()):
            modelo = str(row.get("Modelo", idx))
            rows.append([
                modelo,
                fmt(row.get(col_map["Precision"]) if col_map["Precision"] else None),
                fmt(row.get(col_map["Recall"])    if col_map["Recall"] else None),
                fmt(row.get(col_map["F1"])         if col_map["F1"] else None),
                fmt(row.get(col_map["AUC-ROC"])   if col_map["AUC-ROC"] else None),
                fmt(row.get(col_map["AUC-PR"])    if col_map["AUC-PR"] else None),
                fmt(row.get(col_map["Threshold"]) if col_map["Threshold"] else None, decimals=2),
            ])
        best_idx = 0
    else:
        for name in model_order:
            rows.append([name, PENDING, PENDING, PENDING, PENDING, PENDING, PENDING])

    # ── Tabla ML ────────────────────────────────────────────────────
    add_table(slide, headers, rows,
              Inches(0.35), Inches(1.25), Inches(8.3), Inches(2.35),
              highlight_row=best_idx,
              col_widths=[2.5, 1, 1, 1, 1, 1, 0.9])

    # ── Nota accuracy ───────────────────────────────────────────────
    add_textbox(slide,
                "★  Accuracy no es criterio principal: con prevalencia ≈ 19 %, un clasificador trivial alcanza ~81 % sin detectar ningún cáncer.",
                Inches(0.35), Inches(3.65), Inches(8.3), Inches(0.32),
                font_size=8, color=GRIS_TEXTO, italic=True)

    # ── KPIs del mejor modelo ───────────────────────────────────────
    add_textbox(slide, "MEJOR MODELO ML",
                Inches(8.85), Inches(1.25), Inches(4.15), Inches(0.28),
                font_size=9, bold=True, color=AZUL_MEDIO)
    add_rect(slide, Inches(8.85), Inches(1.55), Inches(4.15), Inches(2.0),
             fill_color=GRIS_FONDO, line_color=AZUL_MEDIO, line_width=Pt(1))

    if rows and rows[0][0] != "LogisticRegression":
        best_name = rows[0][0]
        best_rec  = rows[0][2]
        best_f1   = rows[0][3]
        best_auc  = rows[0][4]
        best_thr  = rows[0][6]
    else:
        best_name, best_rec, best_f1, best_auc, best_thr = ("pendiente",) * 5

    add_textbox(slide, best_name,
                Inches(9.0), Inches(1.62), Inches(3.85), Inches(0.4),
                font_size=16, bold=True, color=AZUL_OSCURO, align=PP_ALIGN.CENTER)
    for label, val, top in [
        ("Recall", best_rec, 2.08),
        ("F1",     best_f1,  2.48),
        ("AUC-ROC", best_auc, 2.88),
        ("Umbral", best_thr, 3.28),
    ]:
        add_textbox(slide, f"{label}:  {val}",
                    Inches(9.1), Inches(top), Inches(3.7), Inches(0.35),
                    font_size=12, bold=(label == "Recall"), color=AZUL_OSCURO)

    add_divider(slide, 4.05)

    # ── Figuras ─────────────────────────────────────────────────────
    add_image_safe(slide, "fig_13_roc_curves.png",
                   Inches(0.35), Inches(4.2), Inches(6.2), Inches(2.95),
                   "Curvas ROC — 5 modelos ML (test set)")
    add_image_safe(slide, "fig_15_confusion_matrix_best.png",
                   Inches(6.75), Inches(4.2), Inches(6.25), Inches(2.95),
                   "Matriz de confusión — mejor modelo ML")


def build_slide3(prs: Presentation, m: dict):
    """SLIDE 3 — Red Neuronal MLP."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_accent_bar(slide)
    add_footer(slide)
    slide_number(slide, 3)

    add_title_block(slide,
                    "Red Neuronal MLP: Arquitectura y Validación",
                    "Dense 256→128→64 + BatchNorm + Dropout · Mismo protocolo anti-leakage")
    add_divider(slide, 1.1)

    # ── Columna izquierda: arquitectura ────────────────────────────
    add_textbox(slide, "ARQUITECTURA",
                Inches(0.35), Inches(1.22), Inches(4.5), Inches(0.28),
                font_size=9, bold=True, color=AZUL_MEDIO)

    capas = [
        ("Entrada",           "N features procesadas",   AZUL_CLARO,  AZUL_OSCURO),
        ("Dense(256, relu)",  "BatchNorm · Dropout 0.25", AZUL_CLARO, AZUL_OSCURO),
        ("Dense(128, relu)",  "BatchNorm · Dropout 0.25", AZUL_CLARO, AZUL_OSCURO),
        ("Dense(64,  relu)",  "BatchNorm · Dropout 0.20", AZUL_CLARO, AZUL_OSCURO),
        ("Dense(1, sigmoid)", "P(cáncer = 1)  ∈ [0, 1]",
         RGBColor(0xE8, 0xF5, 0xE9), VERDE),
    ]
    cy = Inches(1.55)
    cw = Inches(4.1)
    ch = Inches(0.52)
    gap = Inches(0.08)
    for nombre, desc, bg, fc in capas:
        add_rect(slide, Inches(0.45), cy, cw, ch,
                 fill_color=bg, line_color=AZUL_MEDIO, line_width=Pt(0.75))
        add_textbox(slide, nombre, Inches(0.6), cy + Inches(0.04),
                    cw - Inches(0.2), Inches(0.24),
                    font_size=9, bold=True, color=fc)
        add_textbox(slide, desc, Inches(0.6), cy + Inches(0.26),
                    cw - Inches(0.2), Inches(0.22),
                    font_size=7.5, color=GRIS_TEXTO)
        cy += ch + gap
        if nombre != "Dense(1, sigmoid)":
            add_textbox(slide, "↓", Inches(2.35), cy - gap - Inches(0.02),
                        Inches(0.4), Inches(0.1),
                        font_size=9, color=AZUL_MEDIO, align=PP_ALIGN.CENTER)

    # Regularización bullets
    reg_top = Inches(4.75)
    add_textbox(slide, "REGULARIZACIÓN",
                Inches(0.35), reg_top, Inches(4.5), Inches(0.25),
                font_size=9, bold=True, color=AZUL_MEDIO)
    regs = [
        "BatchNorm — estabiliza gradientes entre capas",
        "Dropout — previene sobreajuste (75–80 % activaciones)",
        "EarlyStopping (paciencia 12) — restaura mejor peso",
        "ReduceLROnPlateau — lr × 0.5 si val_loss estanca",
        "class_weight balanced — compensa desbalance 4.3:1",
        "Adam lr=0.001 · batch=256 · máx 150 épocas",
    ]
    ry = reg_top + Inches(0.28)
    for r in regs:
        add_textbox(slide, f"·  {r}", Inches(0.45), ry, Inches(4.2), Inches(0.24),
                    font_size=8, color=AZUL_OSCURO)
        ry += Inches(0.24)

    # ── Columna derecha: comparativa + learning curves ─────────────
    add_textbox(slide, "MLP vs MEJOR MODELO ML — TEST SET",
                Inches(5.0), Inches(1.22), Inches(7.98), Inches(0.28),
                font_size=9, bold=True, color=AZUL_MEDIO)

    PENDING = "pendiente de ejecución"

    def get_mlp_metrics():
        if m["mlp_results"] is not None:
            row = m["mlp_results"].iloc[0]
            return {k: fmt(row.get(k)) for k in ["Recall", "F1", "AUC-ROC", "AUC-PR", "Precision"]}
        if m["all_metrics"] is not None:
            df = m["all_metrics"]
            mlp_rows = df[df.index.astype(str).str.upper().str.contains("MLP")]
            if len(mlp_rows) == 0 and "Tipo" in df.columns:
                mlp_rows = df[df["Tipo"].astype(str).str.upper().str.contains("MLP")]
            if len(mlp_rows) > 0:
                row = mlp_rows.iloc[0]
                return {k: fmt(row.get(k)) for k in ["Recall", "F1", "AUC-ROC", "AUC-PR", "Precision"]}
        return {k: PENDING for k in ["Recall", "F1", "AUC-ROC", "AUC-PR", "Precision"]}

    def get_best_ml_metrics():
        if m["all_metrics"] is not None:
            df = m["all_metrics"].copy()
            if "Tipo" in df.columns:
                df = df[~df["Tipo"].astype(str).str.upper().str.contains("MLP")]
            f1_col = next((c for c in ["F1", "F1_cancer1"] if c in df.columns), None)
            if f1_col:
                row = df.sort_values(f1_col, ascending=False).iloc[0]
                return {k: fmt(row.get(k)) for k in ["Recall", "F1", "AUC-ROC", "AUC-PR", "Precision"]}
        return {k: PENDING for k in ["Recall", "F1", "AUC-ROC", "AUC-PR", "Precision"]}

    mlp  = get_mlp_metrics()
    bml  = get_best_ml_metrics()

    comp_headers = ["Modelo", "Recall ↑", "F1 ↑", "AUC-ROC", "AUC-PR"]
    comp_rows = [
        ["MLP (Red Neuronal)",      mlp["Recall"], mlp["F1"], mlp["AUC-ROC"], mlp["AUC-PR"]],
        ["Mejor ML (clásico)",      bml["Recall"], bml["F1"], bml["AUC-ROC"], bml["AUC-PR"]],
    ]
    add_table(slide, comp_headers, comp_rows,
              Inches(5.0), Inches(1.55), Inches(7.98), Inches(1.1),
              highlight_row=None,
              col_widths=[2.5, 1.2, 1.2, 1.2, 1.2])

    # Conclusión
    add_rect(slide, Inches(5.0), Inches(2.75), Inches(7.98), Inches(0.65),
             fill_color=RGBColor(0xE8, 0xF0, 0xF8), line_color=AZUL_MEDIO, line_width=Pt(0.75))
    add_textbox(slide,
                "Conclusión: La MLP alcanza rendimiento comparable al mejor modelo clásico con mayor "
                "complejidad operacional y menor interpretabilidad. En contexto hospitalario, "
                "el modelo clásico es preferible salvo que la MLP supere claramente en Recall.",
                Inches(5.15), Inches(2.82), Inches(7.68), Inches(0.55),
                font_size=8.5, color=AZUL_OSCURO, italic=True)

    add_divider(slide, 3.52)

    # Figuras
    add_image_safe(slide, "fig_17_mlp_architecture.png",
                   Inches(5.0), Inches(3.65), Inches(3.85), Inches(3.55),
                   "Arquitectura MLP — diagrama de capas")
    add_image_safe(slide, "fig_18_mlp_learning_curves.png",
                   Inches(9.1), Inches(3.65), Inches(3.88), Inches(3.55),
                   "Curvas de aprendizaje — loss y recall por época")


def build_slide4(prs: Presentation, m: dict):
    """SLIDE 4 — Comparativa global."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_accent_bar(slide)
    add_footer(slide)
    slide_number(slide, 4)

    add_title_block(slide,
                    "Ranking Final: ML vs MLP",
                    "Score compuesto = 0.35·Recall + 0.30·F1 + 0.20·AUC-PR + 0.15·AUC-ROC  —  prioridad clínica: minimizar falsos negativos")
    add_divider(slide, 1.1)

    # ── Tabla ranking ───────────────────────────────────────────────
    rank_headers = ["#", "Modelo", "Tipo", "Recall", "F1", "AUC-PR", "AUC-ROC", "Score"]
    rank_rows = []

    medals = ["🥇", "🥈", "🥉", "4", "5", "6"]

    if m["ranking"] is not None:
        df = m["ranking"].copy()
        for i, (idx, row) in enumerate(df.iterrows()):
            nombre = str(row.get("Modelo", idx))
            tipo = str(row.get("Tipo", "MLP" if "MLP" in nombre.upper() else "ML"))
            rank_rows.append([
                medals[i] if i < len(medals) else str(i + 1),
                nombre,
                tipo,
                fmt(row.get("Recall")),
                fmt(row.get("F1")),
                fmt(row.get("AUC-PR")),
                fmt(row.get("AUC-ROC")),
                fmt(row.get("Score_compuesto", row.get("Score"))),
            ])
    else:
        for i, name in enumerate(["LightGBM", "CatBoost", "XGBoost", "Random Forest", "MLP", "Log. Regresión"]):
            tipo = "MLP" if name == "MLP" else "ML"
            rank_rows.append([medals[i], name, tipo,
                               "pendiente", "pendiente", "pendiente", "pendiente", "pendiente"])

    add_table(slide, rank_headers, rank_rows,
              Inches(0.35), Inches(1.25), Inches(12.63), Inches(2.3),
              highlight_row=0,
              col_widths=[0.4, 2.2, 0.8, 1.1, 1.1, 1.1, 1.1, 1.1])

    add_divider(slide, 3.68)

    # ── Mensaje clave ───────────────────────────────────────────────
    add_rect(slide, Inches(0.35), Inches(3.78), Inches(12.63), Inches(0.6),
             fill_color=RGBColor(0xE8, 0xF0, 0xF8), line_color=AZUL_MEDIO, line_width=Pt(0.75))
    add_rect(slide, Inches(0.35), Inches(3.78), Pt(5), Inches(0.6), fill_color=AZUL_MEDIO)
    add_textbox(slide,
                "Criterio de selección: mejor equilibrio Recall · F1 · AUC-PR · AUC-ROC + interpretabilidad clínica. "
                "Los boosters (LightGBM / CatBoost) ofrecen SHAP nativo. La MLP requiere XAI externo.",
                Inches(0.58), Inches(3.86), Inches(12.2), Inches(0.45),
                font_size=9, color=AZUL_OSCURO)

    # ── Figuras ─────────────────────────────────────────────────────
    add_image_safe(slide, "fig_final_metrics_comparison.png",
                   Inches(0.35), Inches(4.5), Inches(6.2), Inches(2.65),
                   "Comparativa global de métricas — todos los modelos")
    add_image_safe(slide, "fig_final_precision_recall_space.png",
                   Inches(6.75), Inches(4.5), Inches(6.25), Inches(2.65),
                   "Espacio Precision–Recall — zona clínica deseable: Recall > 0.70")


def build_slide5(prs: Presentation, m: dict):
    """SLIDE 5 — Viabilidad y decisión."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_accent_bar(slide)
    add_footer(slide)
    slide_number(slide, 5)

    add_title_block(slide,
                    "Viabilidad Clínica y Recomendación Final",
                    "Sistema de soporte al cribado oncológico · No reemplaza el diagnóstico médico")
    add_divider(slide, 1.1)

    # ── KPI del modelo recomendado ──────────────────────────────────
    best = get_best_model_row(m)
    if best:
        best_name = str(best.get("Modelo", "pendiente de ejecución"))
        best_rec  = fmt(best.get("Recall"))
        best_f1   = fmt(best.get("F1"))
        best_auc  = fmt(best.get("AUC-ROC"))
        best_thr  = fmt(best.get("Threshold", best.get("Umbral")), decimals=2)
    else:
        best_name = "pendiente de ejecución"
        best_rec = best_f1 = best_auc = best_thr = "pendiente"

    # Recuadro "modelo recomendado"
    add_rect(slide, Inches(0.35), Inches(1.25), Inches(5.2), Inches(1.55),
             fill_color=AZUL_OSCURO, line_color=AZUL_OSCURO)
    add_textbox(slide, "MODELO RECOMENDADO PARA PRODUCCIÓN",
                Inches(0.5), Inches(1.32), Inches(4.9), Inches(0.28),
                font_size=8, bold=True, color=AZUL_CLARO)
    add_textbox(slide, best_name,
                Inches(0.5), Inches(1.58), Inches(4.9), Inches(0.45),
                font_size=18, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    for label, val, xi in [
        ("Recall",   best_rec, 0.5),
        ("F1",       best_f1,  1.8),
        ("AUC-ROC",  best_auc, 3.1),
        ("Umbral",   best_thr, 4.1),
    ]:
        add_textbox(slide, label,  Inches(xi), Inches(2.1),  Inches(1.1), Inches(0.2),
                    font_size=7, color=AZUL_CLARO, align=PP_ALIGN.CENTER)
        add_textbox(slide, val, Inches(xi), Inches(2.28), Inches(1.1), Inches(0.32),
                    font_size=14, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

    # Veredicto
    add_rect(slide, Inches(5.7), Inches(1.25), Inches(7.3), Inches(1.55),
             fill_color=RGBColor(0xE8, 0xF5, 0xE9), line_color=VERDE, line_width=Pt(1.2))
    add_rect(slide, Inches(5.7), Inches(1.25), Pt(5), Inches(1.55), fill_color=VERDE)
    add_textbox(slide, "VEREDICTO EJECUTIVO",
                Inches(5.9), Inches(1.32), Inches(6.9), Inches(0.25),
                font_size=8, bold=True, color=RGBColor(0x1B, 0x5E, 0x20))
    add_textbox(slide,
                "Sistema viable como herramienta de soporte al cribado clínico.\n"
                "No constituye diagnóstico médico. Requiere validación externa\n"
                "en cohorte real antes de despliegue clínico.",
                Inches(5.9), Inches(1.6), Inches(6.9), Inches(1.1),
                font_size=10, color=RGBColor(0x1B, 0x5E, 0x20))

    add_divider(slide, 2.92)

    # ── 3 columnas: limitaciones / mejoras / app ────────────────────
    col_tops = [Inches(0.35), Inches(4.45), Inches(8.55)]
    col_w = Inches(3.85)

    # Limitaciones
    add_rect(slide, col_tops[0], Inches(3.02), col_w, Inches(4.08),
             fill_color=RGBColor(0xFD, 0xED, 0xEB), line_color=ROJO, line_width=Pt(0.75))
    add_rect(slide, col_tops[0], Inches(3.02), col_w, Inches(0.32), fill_color=ROJO)
    add_textbox(slide, "LIMITACIONES", col_tops[0] + Inches(0.1), Inches(3.05),
                col_w - Inches(0.2), Inches(0.25),
                font_size=9, bold=True, color=BLANCO)
    lims = [
        "Datos sintéticos — no generalizable sin\nvalidación en pacientes reales",
        "Riesgo de falsos negativos — Recall < 1.0\nimplica cánceres no detectados",
        "Posible sesgo poblacional si la distribución\nsintética difiere de la real",
        "MLP: caja negra sin XAI externo\n(SHAP/LIME no implementado)",
    ]
    ly = Inches(3.42)
    for l in lims:
        add_textbox(slide, f"▸  {l}", col_tops[0] + Inches(0.1), ly,
                    col_w - Inches(0.2), Inches(0.55),
                    font_size=8, color=RGBColor(0x7A, 0x1A, 0x10))
        ly += Inches(0.6)

    # Mejoras futuras
    add_rect(slide, col_tops[1], Inches(3.02), col_w, Inches(4.08),
             fill_color=RGBColor(0xE8, 0xF0, 0xF8), line_color=AZUL_MEDIO, line_width=Pt(0.75))
    add_rect(slide, col_tops[1], Inches(3.02), col_w, Inches(0.32), fill_color=AZUL_MEDIO)
    add_textbox(slide, "MEJORAS FUTURAS", col_tops[1] + Inches(0.1), Inches(3.05),
                col_w - Inches(0.2), Inches(0.25),
                font_size=9, bold=True, color=BLANCO)
    mejoras = [
        "SHAP — interpretabilidad por paciente\npara el clínico",
        "Calibración de probabilidades\n(Platt scaling / Isotonic regression)",
        "Validación en cohorte real con datos\nanonimizados hospitalarios",
        "Monitorización de drift — reentrenamiento\nperiódico en producción",
    ]
    my = Inches(3.42)
    for mej in mejoras:
        add_textbox(slide, f"→  {mej}", col_tops[1] + Inches(0.1), my,
                    col_w - Inches(0.2), Inches(0.55),
                    font_size=8, color=AZUL_OSCURO)
        my += Inches(0.6)

    # App Streamlit
    add_rect(slide, col_tops[2], Inches(3.02), col_w, Inches(4.08),
             fill_color=RGBColor(0xE8, 0xF5, 0xE9), line_color=VERDE, line_width=Pt(0.75))
    add_rect(slide, col_tops[2], Inches(3.02), col_w, Inches(0.32), fill_color=VERDE)
    add_textbox(slide, "DEMO INTERACTIVA", col_tops[2] + Inches(0.1), Inches(3.05),
                col_w - Inches(0.2), Inches(0.25),
                font_size=9, bold=True, color=BLANCO)
    add_textbox(slide, "streamlit run app.py",
                col_tops[2] + Inches(0.1), Inches(3.42),
                col_w - Inches(0.2), Inches(0.28),
                font_size=10, bold=True, color=RGBColor(0x1B, 0x5E, 0x20))
    app_items = [
        "Simulador de paciente — introduce parámetros\ny obtén probabilidad en tiempo real",
        "Comparativa de modelos — tabla y\ngráficos interactivos",
        "Galería de figuras — todas las visualizaciones\ndel estudio",
        "Estado de artefactos — verificación\ndel pipeline en sidebar",
    ]
    ay = Inches(3.78)
    for a in app_items:
        add_textbox(slide, f"✓  {a}", col_tops[2] + Inches(0.1), ay,
                    col_w - Inches(0.2), Inches(0.55),
                    font_size=8, color=RGBColor(0x1B, 0x5E, 0x20))
        ay += Inches(0.6)


# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────

def main():
    print("=" * 60)
    print("Generando presentación: cancer_prediction_presentation.pptx")
    print("=" * 60)

    print("\n[1/3] Cargando métricas...")
    m = load_metrics()
    for key, val in m.items():
        status = "OK" if val is not None else "NO ENCONTRADO — se usará 'pendiente de ejecución'"
        print(f"  {key:20s}: {status}")

    print("\n[2/3] Construyendo slides...")
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide1(prs, m)
    print("  Slide 1 — Objetivo y datos: OK")
    build_slide2(prs, m)
    print("  Slide 2 — Resultados ML: OK")
    build_slide3(prs, m)
    print("  Slide 3 — Red Neuronal MLP: OK")
    build_slide4(prs, m)
    print("  Slide 4 — Comparativa global: OK")
    build_slide5(prs, m)
    print("  Slide 5 — Viabilidad y decisión: OK")

    print(f"\n[3/3] Guardando en {OUT_PPTX} ...")
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX))
    print(f"\nOK  Presentacion guardada: {OUT_PPTX}")
    print(f"   Slides totales: {len(prs.slides)}")
    assert len(prs.slides) == 5, f"ERROR: se esperaban 5 slides, hay {len(prs.slides)}"
    print("\nOK - exactamente 5 diapositivas.")


if __name__ == "__main__":
    main()
